import os
import subprocess
import logging
from backend.converters.base import BaseConverter
from backend.utils.sys_info import get_libreoffice_path

logger = logging.getLogger("slides_converter")

class SlidesConverter(BaseConverter):
    def can_convert(self, from_ext: str, to_ext: str) -> bool:
        from_ext = from_ext.lower().strip('.')
        to_ext = to_ext.lower().strip('.')
        
        if from_ext in ["pptx", "ppt", "odp"]:
            return to_ext in ["pdf", "png", "jpg"]
        return False

    def is_available(self) -> bool:
        return True

    def convert(self, input_path: str, output_path: str, options: dict = None) -> bool:
        from_ext = os.path.splitext(input_path)[1].lower().strip('.')
        to_ext = os.path.splitext(output_path)[1].lower().strip('.')

        # Check for LibreOffice (best quality)
        lo_path = get_libreoffice_path()
        if lo_path:
            return self._convert_with_libreoffice(lo_path, input_path, output_path, to_ext)

        # Fallback presentation to PDF if LibreOffice is missing
        if from_ext == "pptx" and to_ext == "pdf":
            logger.warning("LibreOffice not detected. Using python-pptx fallback.")
            return self._pptx_to_pdf_fallback(input_path, output_path)

        raise RuntimeError(
            f"Powerpoint conversion from .{from_ext} to .{to_ext} without LibreOffice is not supported. "
            f"Please install LibreOffice."
        )

    def _convert_with_libreoffice(self, lo_path: str, input_path: str, output_path: str, to_ext: str) -> bool:
        logger.info(f"Converting slides using LibreOffice: {input_path} -> {to_ext}")
        out_dir = os.path.dirname(output_path)
        
        # If outputting to images (png/jpg), LibreOffice can export them
        lo_to_ext = to_ext
        if to_ext in ["png", "jpg"]:
            lo_to_ext = "pdf" # Convert to PDF first, then we can extract images using PyMuPDF later in the pipeline
        
        cmd = [
            lo_path,
            "--headless",
            "--convert-to",
            f"{lo_to_ext}:impress_pdf_Export",
            "--outdir",
            out_dir,
            input_path
        ]
        
        try:
            result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=30)
            if result.returncode != 0:
                logger.error(f"LibreOffice failed: {result.stderr}")
                raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

            in_basename = os.path.basename(input_path)
            in_name_no_ext = os.path.splitext(in_basename)[0]
            lo_output_name = f"{in_name_no_ext}.{lo_to_ext}"
            lo_output_path = os.path.join(out_dir, lo_output_name)

            if to_ext in ["png", "jpg"]:
                # Convert the PDF pages to PNG/JPG using PyMuPDF (fitz)
                import fitz
                doc = fitz.open(lo_output_path)
                # Save first page as slide thumbnail or multiple pages as requested
                # By default save page 0
                page = doc.load_page(0)
                pix = page.get_pixmap()
                pix.save(output_path)
                doc.close()
                if os.path.exists(lo_output_path):
                    os.remove(lo_output_path)
                return True
            
            if os.path.exists(lo_output_path) and os.path.abspath(lo_output_path) != os.path.abspath(output_path):
                if os.path.exists(output_path):
                    os.remove(output_path)
                os.rename(lo_output_path, output_path)
                return True
            elif os.path.exists(output_path):
                return True
            else:
                raise FileNotFoundError(f"LibreOffice output not found at {lo_output_path}")
        except subprocess.TimeoutExpired:
            logger.error("LibreOffice presentation conversion timed out.")
            raise RuntimeError("LibreOffice presentation conversion timed out.")

    def _pptx_to_pdf_fallback(self, input_path: str, output_path: str) -> bool:
        from pptx import Presentation
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors

        prs = Presentation(input_path)
        pdf = SimpleDocTemplate(output_path, pagesize=landscape(letter), leftMargin=36, rightMargin=36, topMargin=36, bottomMargin=36)
        styles = getSampleStyleSheet()
        
        title_style = ParagraphStyle(
            'SlidesTitle',
            parent=styles['Heading1'],
            fontSize=24,
            leading=28,
            textColor=colors.HexColor('#1e3a8a'),
            spaceAfter=15
        )
        
        body_style = ParagraphStyle(
            'SlidesBody',
            parent=styles['Normal'],
            fontSize=14,
            leading=18,
            spaceAfter=10
        )

        story = []
        for i, slide in enumerate(prs.slides):
            if i > 0:
                story.append(PageBreak())
                
            story.append(Paragraph(f"Slide {i+1}", styles['Normal']))
            story.append(Spacer(1, 10))

            # Extract text from shapes
            has_content = False
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue
                    has_content = True
                    # Let's simple check if it looks like a slide title
                    if shape.name.lower().startswith("title") or shape == slide.shapes[0]:
                        story.append(Paragraph(text, title_style))
                    else:
                        story.append(Paragraph(text, body_style))
            
            if not has_content:
                story.append(Paragraph("[Empty Slide or Image-Only Slide]", body_style))

        pdf.build(story)
        return True
